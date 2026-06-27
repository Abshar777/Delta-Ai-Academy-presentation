#!/usr/bin/env python3
"""Option B: one slide per scene, each embedding that scene's rendered MP4,
set to auto-play on entry and auto-advance when the clip ends.
Reads per-scene clips from a manifest produced by render_clips.js."""
import json
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from lxml import etree

CX, CY = 12191695, 6858000

DECKS = [
    ("/tmp/clips_navy", "output/Delta-AI-Academy-Keynote-Dark-Animated.pptx", "06121F"),
    ("/tmp/clips_oled", "output/Delta-AI-Academy-Keynote-Dark-OLED-Animated.pptx", "070707"),
]

def load_manifest(clips_dir):
    with open(f"{clips_dir}/manifest.json") as f:
        data = json.load(f)
    return [(d["clip"], d["poster"], d["dur_ms"]) for d in data]

def set_bg(slide, hexcolor):
    cSld = slide._element.find(qn('p:cSld'))
    bg = etree.Element(qn('p:bg'))
    bgPr = etree.SubElement(bg, qn('p:bgPr'))
    fill = etree.SubElement(bgPr, qn('a:solidFill'))
    etree.SubElement(fill, qn('a:srgbClr')).set('val', hexcolor)
    etree.SubElement(bgPr, qn('a:effectLst'))
    cSld.insert(0, bg)

def add_transition(slide, adv_ms):
    t = etree.SubElement(slide._element, qn('p:transition'))
    t.set('spd', 'med'); t.set('advClick', '1'); t.set('advTm', str(adv_ms))
    etree.SubElement(t, qn('p:fade'))

TIMING_TMPL = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
 <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
  <p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
   <p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst>
    <p:par><p:cTn id="4" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
     <p:par><p:cTn id="5" presetClass="mediacall" presetID="0" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
      <p:cmd type="call" cmd="playFrom(0.0)"><p:cBhvr><p:cTn id="6" dur="{dur}" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:cmd>
     </p:childTnLst></p:cTn></p:par>
    </p:childTnLst></p:cTn></p:par>
   </p:childTnLst></p:cTn></p:par>
  </p:childTnLst></p:cTn>
  <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
  <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
  </p:seq>
 </p:childTnLst></p:cTn></p:par></p:tnLst>
</p:timing>"""

def add_autoplay(slide, spid, dur_ms):
    xml = TIMING_TMPL.format(dur=dur_ms, spid=spid)
    slide._element.append(etree.fromstring(xml))

def build(clips_dir, out_path, bg_hex):
    meta = load_manifest(clips_dir)
    prs = Presentation()
    prs.slide_width = Emu(CX); prs.slide_height = Emu(CY)
    blank = prs.slide_layouts[6]
    for clip, poster, dur_ms in meta:
        slide = prs.slides.add_slide(blank)
        set_bg(slide, bg_hex)
        mv = slide.shapes.add_movie(clip, 0, 0, Emu(CX), Emu(CY),
                                    poster_frame_image=poster, mime_type='video/mp4')
        add_autoplay(slide, mv.shape_id, dur_ms)
        add_transition(slide, dur_ms + 250)
    prs.save(out_path)
    total = sum(m[2] for m in meta) / 1000.0
    print(f"saved {out_path}  ({len(meta)} slides, ~{total:.0f}s)")

if __name__ == "__main__":
    for cd, out, bg in DECKS:
        build(cd, out, bg)
