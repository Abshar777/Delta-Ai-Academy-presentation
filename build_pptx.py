#!/usr/bin/env python3
"""Assemble a dark-deck PPTX: one full-bleed slide per scene + Fade transitions."""
import sys, copy
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from lxml import etree

# slide-ordered timestamps used to name the rendered spot frames (scene order)
TS = ["5.6","13.1","20.1","26.6","33.6","41.6","50.1","59.6",
      "72.1","77.1","91.1","98.1","109.1","120.1","132.1","140.1"]

CX, CY = 12191695, 6858000   # match existing keynote (16:9)

def fname(t):
    return "spot_" + t.replace(".", "_") + ".jpg"

def add_fade(slide, spd="med"):
    # <p:transition spd="med"><p:fade/></p:transition>
    t = etree.SubElement(slide._element, qn('p:transition'))
    t.set('spd', spd)
    etree.SubElement(t, qn('p:fade'))

def set_bg(slide, hexcolor):
    cSld = slide._element.find(qn('p:cSld'))
    bg = etree.Element(qn('p:bg'))
    bgPr = etree.SubElement(bg, qn('p:bgPr'))
    fill = etree.SubElement(bgPr, qn('a:solidFill'))
    clr = etree.SubElement(fill, qn('a:srgbClr'))
    clr.set('val', hexcolor)
    etree.SubElement(bgPr, qn('a:effectLst'))
    cSld.insert(0, bg)

def build(frames_dir, out_path, bg_hex):
    prs = Presentation()
    prs.slide_width = Emu(CX)
    prs.slide_height = Emu(CY)
    blank = prs.slide_layouts[6]
    for t in TS:
        slide = prs.slides.add_slide(blank)
        set_bg(slide, bg_hex)
        slide.shapes.add_picture(f"{frames_dir}/{fname(t)}", 0, 0, width=Emu(CX), height=Emu(CY))
        add_fade(slide)
    prs.save(out_path)
    print(f"saved {out_path}  ({len(TS)} slides)")

if __name__ == "__main__":
    build("/tmp/sc_navy",
          "output/Delta-AI-Academy-Keynote-Dark.pptx", "06121F")
    build("/tmp/sc_oled",
          "output/Delta-AI-Academy-Keynote-Dark-OLED.pptx", "070707")
